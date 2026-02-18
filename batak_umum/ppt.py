# ppt.py

from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement


def format_judul_acara(text):
    if text is None:
        return ""
    text = str(text).strip()

    if "TANGIANG PANGUJUNGI" in text.upper():
        return "TANGIANG PANGUJUNGI"
    if text.endswith(":"):
        text = text[:-1].strip()

    text_upper = text.upper().replace(" ", "")

    if "MARENDE" in text_upper:
        pattern = r"(MARENDE)\s+(.*?)\s+([“\"].*[”\"])"
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            return f"{match.group(1)}\n{match.group(2)}\n{match.group(3)}"

    if "KOOR" in text_upper and "-" in text:
        parts = [p.strip() for p in text.split("-") if p.strip()]
        prefix_match = re.match(r"^([^:]+)", text)
        prefix = prefix_match.group(1).strip() if prefix_match else "K O O R"
        formatted_parts = []
        for p in parts:
            clean_p = re.sub(rf"^{re.escape(prefix)}", "",
                             p, flags=re.IGNORECASE).strip()
            if clean_p:
                formatted_parts.append(f"{prefix} - {clean_p}")
        return "\n".join(formatted_parts)

    return text


def apply_radical_styling(paragraph, font_size=60, alignment=PP_ALIGN.CENTER, font_color=RGBColor(0, 0, 0)):
    paragraph.alignment = alignment
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ['buNone', 'buAutoNum', 'buChar', 'buBlip']:
        child = pPr.find(qn(f'a:{tag}'))
        if child is not None:
            pPr.remove(child)

    buNone = OxmlElement('a:buNone')
    pPr.insert(0, buNone)
    pPr.set('marL', '0')
    pPr.set('indent', '0')

    for run in paragraph.runs:
        run.font.name = 'Verdana'
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color


def create_styled_slide(prs, content_text, font_size=60, alignment=PP_ALIGN.CENTER, header_text=None):
    if not str(content_text).strip():
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if header_text:
        h_width = prs.slide_width - Inches(1)
        h_height = Inches(0.5)
        h_left = (prs.slide_width - h_width) // 2
        h_top = Inches(0.2)

        header_box = slide.shapes.add_textbox(h_left, h_top, h_width, h_height)
        h_p = header_box.text_frame.paragraphs[0]
        h_p.text = str(header_text).replace("\n", " ")
        apply_radical_styling(
            h_p, font_size=18, alignment=PP_ALIGN.CENTER, font_color=RGBColor(255, 165, 0))

    width = prs.slide_width - Inches(0.5)
    height = prs.slide_height - Inches(1.5)
    left = (prs.slide_width - width) // 2
    top = (prs.slide_height - height) // 2

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = str(content_text)
    apply_radical_styling(p, font_size=font_size, alignment=alignment)


def create_cover_slide(prs, minggu, topik, tanggal):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    width_main = prs.slide_width - Inches(1)
    height_main = Inches(4)
    left_main = (prs.slide_width - width_main) // 2
    top_main = Inches(1.5)

    tb_main = slide.shapes.add_textbox(
        left_main, top_main, width_main, height_main)
    tf_main = tb_main.text_frame
    tf_main.word_wrap = True
    tf_main.vertical_anchor = MSO_ANCHOR.MIDDLE

    p_main = tf_main.paragraphs[0]
    p_main.text = f"{minggu}\n\"{topik}\""
    apply_radical_styling(p_main, font_size=54)

    width_date = prs.slide_width - Inches(1)
    height_date = Inches(1)
    left_date = (prs.slide_width - width_date) // 2
    top_date = prs.slide_height - Inches(2)

    tb_date = slide.shapes.add_textbox(
        left_date, top_date, width_date, height_date)
    p_date = tb_date.text_frame.paragraphs[0]
    p_date.text = tanggal
    apply_radical_styling(p_date, font_size=32)


def generate_slides(prs, cover_data, sections):
    if cover_data and not cover_data.get('skip_cover', False):
        create_cover_slide(
            prs, cover_data['minggu'], cover_data['topik'], cover_data['tanggal'])

    for section in sections:
        raw_judul = section.get('judul', '').strip()
        isi_raw = section.get('isi', [])

        if not raw_judul and not isi_raw:
            continue

        formatted_judul = format_judul_acara(raw_judul)
        create_styled_slide(prs, formatted_judul)

        if isi_raw:
            raw_judul_upper = raw_judul.upper().replace(" ", "")
            is_lagu = "MARENDE" in raw_judul_upper
            is_epistel = "EPISTEL" in raw_judul_upper or "E P I S T E L" in raw_judul_upper

            if is_lagu:
                for line in isi_raw:
                    clean_line = " ".join(str(line).split())
                    if clean_line and not any(clean_line.startswith(x) for x in ["[", "---"]):
                        create_styled_slide(
                            prs,
                            clean_line,
                            font_size=60,
                            alignment=PP_ALIGN.CENTER,
                            header_text=formatted_judul
                        )
            elif is_epistel:
                for line in isi_raw:
                    clean_line = str(line).strip()
                    if clean_line and not any(clean_line.startswith(x) for x in ["[", "---"]):
                        create_styled_slide(
                            prs, clean_line, font_size=36, alignment=PP_ALIGN.LEFT)
