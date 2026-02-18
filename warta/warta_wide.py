# warta/warta_wide.py

import io
import re
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def add_title_slide(prs, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    txBox = slide.shapes.add_textbox(Pt(50), Pt(
        0), prs.slide_width - Pt(100), prs.slide_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text.upper()
    run.font.name = 'Verdana'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 0, 0)


def add_content_slide(prs, content, force_bold=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    txBox = slide.shapes.add_textbox(Pt(50), Pt(
        50), prs.slide_width - Pt(100), prs.slide_height - Pt(100))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = content.strip()
    run.font.name = 'Verdana'
    run.font.bold = force_bold
    run.font.size = Pt(44)
    run.font.color.rgb = RGBColor(0, 0, 0)


def extract_images(doc, prs):
    for rel_id, part in doc.part.related_parts.items():
        if "image" in part.content_type:
            image_stream = io.BytesIO(part.blob)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            margin = Inches(0.5)
            pic = slide.shapes.add_picture(image_stream, margin, margin)
            ratio = min((prs.slide_width - 2*margin) / pic.width,
                        (prs.slide_height - 2*margin) / pic.height)
            pic.width = int(pic.width * ratio)
            pic.height = int(pic.height * ratio)
            pic.left = int((prs.slide_width - pic.width) / 2)
            pic.top = int((prs.slide_height - pic.height) / 2)


def generate_warta(doc, prs):
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

    all_paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    pre_topic_content = []
    main_sections = []

    found_topic = False
    current_section_title = None
    current_section_sentences = []

    for text in all_paragraphs:
        is_heading = bool(re.match(r'^\d+\.', text))

        if "TOPIK" in text.upper() and not is_heading:
            found_topic = True
            continue

        if not found_topic:
            if not text.startswith('=='):
                pre_topic_content.append(text)
            continue

        if is_heading:
            if current_section_title or current_section_sentences:
                main_sections.append(
                    (current_section_title, current_section_sentences))

            current_section_title = text
            current_section_sentences = []
        else:
            if (current_section_title or main_sections) and not text.startswith('=='):
                sentences = re.split(r'(?<=[.!?])\s+', text)
                current_section_sentences.extend(
                    [s.strip() for s in sentences if s.strip()])

    if current_section_title or current_section_sentences:
        main_sections.append(
            (current_section_title, current_section_sentences))

    if pre_topic_content:
        add_content_slide(prs, "\n".join(pre_topic_content), force_bold=True)

    for title, sentences in main_sections:
        if title:
            add_title_slide(prs, title)

        accumulator = []
        for sentence in sentences:
            if len(" ".join(accumulator + [sentence]).split()) > 35:
                if accumulator:
                    add_content_slide(prs, " ".join(accumulator))
                    accumulator = [sentence]
                else:
                    add_content_slide(prs, sentence)
                    accumulator = []
            else:
                accumulator.append(sentence)

        if accumulator:
            add_content_slide(prs, " ".join(accumulator))

    extract_images(doc, prs)
