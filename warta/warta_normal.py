# warta/warta_normal.py

import io
import re
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


def add_title_slide(prs, title_text, is_romawi=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    txBox = slide.shapes.add_textbox(Pt(50), Pt(
        0), prs.slide_width - Pt(100), prs.slide_height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text.upper()
    run.font.name = 'Verdana'
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 0, 0)


def add_content_slide(prs, content, force_bold=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    txBox = slide.shapes.add_textbox(Pt(50), Pt(
        50), prs.slide_width - Pt(100), prs.slide_height - Pt(100))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = content.strip()
    run.font.name = 'Verdana'
    run.font.bold = force_bold
    run.font.size = Pt(48)
    run.font.color.rgb = RGBColor(0, 0, 0)


def extract_images(doc, prs):
    for rel_id, part in doc.part.related_parts.items():
        if "image" in part.content_type:
            image_stream = io.BytesIO(part.blob)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            margin = Inches(0.5)
            pic = slide.shapes.add_picture(image_stream, margin, margin)
            ratio = min((prs.slide_width - 2*margin) / pic.width,
                        (prs.slide_height - 2*margin) / pic.height)
            pic.width = int(pic.width * ratio)
            pic.height = int(pic.height * ratio)
            pic.left = int((prs.slide_width - pic.width) / 2)
            pic.top = int((prs.slide_height - pic.height) / 2)


def generate_warta(doc, prs):
    prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)

    all_paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    pre_heading_content = []
    main_sections = []
    found_any_heading = False
    current_section_title = None
    current_section_words = []

    re_romawi = r'^[IVXLC]+\.\s'
    re_normal = r'^\d+\.\s'
    exclude_pattern = r'(?i)pelayan\s+minggu\s+ini'

    for text in all_paragraphs:
        if text.startswith('==') or re.search(exclude_pattern, text):
            continue

        is_romawi = bool(re.match(re_romawi, text))
        is_normal = bool(re.match(re_normal, text))

        if is_romawi or is_normal:
            if not found_any_heading:
                if pre_heading_content:
                    add_content_slide(prs, "\n".join(
                        pre_heading_content), force_bold=True)
                found_any_heading = True

            if current_section_title or current_section_words:
                main_sections.append(
                    (current_section_title, current_section_words))

            current_section_title = text
            current_section_words = []
        else:
            if not found_any_heading:
                pre_heading_content.append(text)
            else:
                words = text.split()
                current_section_words.extend(words)

    if not found_any_heading and pre_heading_content:
        add_content_slide(prs, "\n".join(pre_heading_content), force_bold=True)

    if current_section_title or current_section_words:
        main_sections.append((current_section_title, current_section_words))

    for title, words in main_sections:
        if title:
            is_rom = bool(re.match(re_romawi, title))
            add_title_slide(prs, title, is_romawi=is_rom)

        if words:
            for i in range(0, len(words), 20):
                chunk = words[i:i + 20]
                add_content_slide(prs, " ".join(chunk))

    extract_images(doc, prs)
