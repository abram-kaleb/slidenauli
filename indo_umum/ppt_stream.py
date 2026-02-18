import re
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

# generate_slides.py


def format_judul_acara(text):
    if text is None:
        return ""
    text = str(text).strip()
    if "DOA PENUTUP" in text.upper():
        return "DOA PENUTUP"
    if text.endswith(":"):
        text = text[:-1].strip()

    text_upper = text.upper().replace(" ", "")

    if any(keyword in text_upper for keyword in ["KJNO", "KJ.", "KJ", "BE", "PKJ", "NKB", "BN.HKBP", "BNHKBP"]):
        pattern = r"(.*?)\s+((?:KJ|BE|PKJ|NKB|BN\.?\s*HKBP)\.?\s*\d+.*)"
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            bagian_awal = match.group(1).strip()
            detail_lagu = match.group(2).strip()

            detail_lagu = re.split(
                r"\s+BL\.", detail_lagu, flags=re.IGNORECASE)[0].strip()
            detail_lagu = re.sub(r"\s+", " ", detail_lagu)

            if not bagian_awal:
                bagian_awal = "BERNYANYI"
            return f"{bagian_awal}\n{detail_lagu}"

    if "KOOR" in text_upper and "-" in text:
        parts = [p.strip() for p in text.split("-") if p.strip()]
        prefix_match = re.match(r"^([^:]+)", text)
        prefix = prefix_match.group(1).strip() if prefix_match else "K O O R"
        formatted_parts = []
        for p in parts:
            clean_p = re.sub(rf"^{re.escape(prefix)}", "",
                             p, flags=re.IGNORECASE).strip()
            if clean_p:
                formatted_parts.append(f"{prefix} - {clean_p}")
        return "\n".join(formatted_parts)

    return text


def apply_radical_styling(paragraph, font_size=40, alignment=PP_ALIGN.CENTER, font_color=RGBColor(255, 255, 255)):
    paragraph.alignment = alignment
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ['buNone', 'buAutoNum', 'buChar', 'buBlip']:
        child = pPr.find(qn(f'a:{tag}'))
        if child is not None:
            pPr.remove(child)

    buNone = OxmlElement('a:buNone')
    pPr.insert(0, buNone)
    pPr.set('marL', '0')
    pPr.set('indent', '0')

    for run in paragraph.runs:
        run.font.name = 'Amasis MT Pro Black'
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color


def create_styled_slide(prs, content_text, font_size=40, alignment=PP_ALIGN.CENTER, header_text=None, mode="Projector"):
    if not str(content_text).strip():
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if mode == "YouTube":
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        rect_height = Inches(1.8)

        bg_shape = slide.shapes.add_shape(
            1, Inches(0), prs.slide_height - rect_height,
            prs.slide_width, rect_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(0, 255, 0)
        bg_shape.line.fill.background()

        width = prs.slide_width - Inches(1.5)
        height = rect_height - Inches(0.4)
        left = (prs.slide_width - width) // 2
        top = prs.slide_height - rect_height + Inches(0.2)

        txt_bg = slide.shapes.add_shape(
            1, left, top, width, height
        )
        txt_bg.fill.solid()
        txt_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
        txt_bg.line.fill.background()

        current_font_color = RGBColor(255, 255, 255)
    else:
        width = prs.slide_width - Inches(0.5)
        height = prs.slide_height - Inches(1.5)
        left = (prs.slide_width - width) // 2
        top = (prs.slide_height - height) // 2
        current_font_color = RGBColor(0, 0, 0)

        if header_text:
            h_width = prs.slide_width - Inches(1)
            h_height = Inches(0.5)
            h_left = (prs.slide_width - h_width) // 2
            h_top = Inches(0.2)
            header_box = slide.shapes.add_textbox(
                h_left, h_top, h_width, h_height)
            h_p = header_box.text_frame.paragraphs[0]
            h_p.text = str(header_text).replace("\n", " ")
            apply_radical_styling(
                h_p, font_size=18, alignment=PP_ALIGN.CENTER, font_color=RGBColor(255, 165, 0))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = str(content_text)
    apply_radical_styling(p, font_size=font_size,
                          alignment=alignment, font_color=current_font_color)


def create_cover_slide(prs, minggu, topik, tanggal, mode="Projector"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if mode == "YouTube":
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        rect_height = Inches(1.8)

        bg_shape = slide.shapes.add_shape(
            1, Inches(0), prs.slide_height - rect_height,
            prs.slide_width, rect_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(0, 255, 0)
        bg_shape.line.fill.background()

        width = prs.slide_width - Inches(1.5)
        height = rect_height - Inches(0.4)
        left = (prs.slide_width - width) // 2
        top = prs.slide_height - rect_height + Inches(0.2)

        txt_bg = slide.shapes.add_shape(
            1, left, top, width, height
        )
        txt_bg.fill.solid()
        txt_bg.fill.fore_color.rgb = RGBColor(0, 0, 0)
        txt_bg.line.fill.background()

        tb_main = slide.shapes.add_textbox(left, top, width, height)
        tf_main = tb_main.text_frame
        tf_main.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_main = tf_main.paragraphs[0]
        p_main.text = str(minggu).upper()
        apply_radical_styling(p_main, font_size=40,
                              font_color=RGBColor(255, 255, 255))
    else:
        width_main = prs.slide_width - Inches(1)
        height_main = Inches(4)
        left_main = (prs.slide_width - width_main) // 2
        top_main = Inches(1.5)
        tb_main = slide.shapes.add_textbox(
            left_main, top_main, width_main, height_main)
        tf_main = tb_main.text_frame
        tf_main.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_main = tf_main.paragraphs[0]
        p_main.text = f"{minggu}\n\"{topik}\""
        apply_radical_styling(p_main, font_size=54,
                              font_color=RGBColor(0, 0, 0))

        width_date = prs.slide_width - Inches(1)
        height_date = Inches(1)
        left_date = (prs.slide_width - width_date) // 2
        top_date = prs.slide_height - Inches(2)
        tb_date = slide.shapes.add_textbox(
            left_date, top_date, width_date, height_date)
        p_date = tb_date.text_frame.paragraphs[0]
        p_date.text = tanggal
        apply_radical_styling(p_date, font_size=32,
                              font_color=RGBColor(0, 0, 0))

# generate_slides.py


def generate_slides(prs, cover_data, sections):
    mode = cover_data.get('mode', 'Projector')
    if mode == "YouTube":
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    if cover_data and not cover_data.get('skip_cover', False):
        create_cover_slide(
            prs, cover_data['minggu'], cover_data['topik'], cover_data['tanggal'], mode=mode)
    for section in sections:
        raw_judul = section.get('judul', '').strip()
        isi_raw = section.get('isi', [])
        if not raw_judul and not isi_raw:
            continue
        formatted_judul = format_judul_acara(raw_judul)
        create_styled_slide(prs, formatted_judul, mode=mode)

        raw_judul_upper = raw_judul.upper().replace(" ", "")

        skip_content_keywords = [
            "VOTUM",
            "HUKUM",
            "DOAPENGAKUAN",
            "PENGAKUANIMAN",
            "DOAPENUTUP"
        ]

        if any(kw in raw_judul_upper for kw in skip_content_keywords):
            continue

        if isi_raw:
            is_lagu = any(kw in raw_judul_upper for kw in [
                          "KJ", "BE", "PKJ", "NKB", "BERNYANYI"])
            if is_lagu:
                for line in isi_raw:
                    clean_line = " ".join(str(line).split())
                    if clean_line:
                        create_styled_slide(
                            prs, clean_line, font_size=40, alignment=PP_ALIGN.CENTER, header_text=formatted_judul, mode=mode)
            elif "EPISTEL" in raw_judul_upper or "KHOTBAH" in raw_judul_upper:
                for line in isi_raw:
                    clean_line = str(line).strip()
                    if clean_line:
                        create_styled_slide(
                            prs, clean_line, font_size=32, alignment=PP_ALIGN.LEFT, mode=mode)
            elif "WARTA" in raw_judul_upper:
                for line in isi_raw:
                    clean_line = str(line).strip()
                    if clean_line:
                        create_styled_slide(
                            prs, clean_line, font_size=28, alignment=PP_ALIGN.LEFT, mode=mode)
            else:
                for line in isi_raw:
                    clean_line = str(line).strip()
                    if clean_line:
                        create_styled_slide(
                            prs, clean_line, font_size=36, alignment=PP_ALIGN.CENTER, mode=mode)
