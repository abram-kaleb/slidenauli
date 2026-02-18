# logic.py

import re
import os
import random
from io import BytesIO
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import warta.warta_normal as warta_normal
import warta.warta_wide as warta_wide


def detect_format(doc):
    full_text = "\n".join([p.text for p in doc.paragraphs[:40]]).lower()
    if "warta" in full_text:
        if any(x in full_text for x in ["remaja", "naposobulung"]):
            return "Warta Remaja"
        return "Warta Jemaat"
    elif "sekolah minggu" in full_text or "skm" in full_text:
        return "Sekolah Minggu (SKM)"
    elif any(x in full_text for x in ["sore", "pukul 17", "pukul 18"]):
        return "Ibadah Sore"
    elif any(x in full_text for x in ["remaja", "naposobulung"]):
        return "Ibadah Remaja"
    elif any(x in full_text for x in ["agenda", "parmingguon", "pukul 07", "pukul 09"]):
        return "Ibadah Batak Umum"
    elif any(x in full_text for x in ["tata ibadah", "pukul 10"]):
        return "Ibadah Indonesia Umum"
    return "Unknown"


def apply_background(prs, slide, bg_path):
    if bg_path and os.path.exists(bg_path):
        pic = slide.shapes.add_picture(
            bg_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, width=prs.slide_width, height=prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)

        spPr = overlay.element.xpath('.//p:spPr')[0]
        solidFill = spPr.xpath('.//a:solidFill')[0]
        srgbClr = solidFill.xpath('.//a:srgbClr')[0]
        etree.SubElement(
            srgbClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha', val='60000')

        overlay.line.fill.background()
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        slide.shapes._spTree.remove(overlay._element)
        slide.shapes._spTree.insert(3, overlay._element)


def set_font_white(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255)


def merge_and_generate(warta_doc, cover_info, data_isi, gen_slides_func, warta_mode):
    prs = Presentation()
    use_bg = cover_info.get("use_bg", False)
    bg_files = []

    if use_bg:
        bg_folder = "pics"
        if os.path.exists(bg_folder):
            bg_files = [os.path.join(bg_folder, f) for f in os.listdir(
                bg_folder) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]

    target_kws = ["WARTA", "TINGTING", "TING TING", "TING-TING"]

    gen_slides_func(prs, cover_info, [])

    if use_bg and bg_files:
        bg_main = random.choice(bg_files)
        apply_background(prs, prs.slides[0], bg_main)
        set_font_white(prs.slides[0])

    for section in data_isi:
        judul = section.get('judul', '').upper()
        is_warta_section = any(kw in judul for kw in target_kws)

        start_idx = len(prs.slides)

        cover_info['skip_cover'] = True
        gen_slides_func(prs, cover_info, [section])

        end_idx = len(prs.slides)

        bg_to_use = random.choice(bg_files) if use_bg and bg_files else None

        for i in range(start_idx, end_idx):
            if bg_to_use:
                apply_background(prs, prs.slides[i], bg_to_use)
                set_font_white(prs.slides[i])

        if is_warta_section and warta_doc:
            if warta_mode == "Normal":
                warta_normal.generate_warta(warta_doc, prs)
            else:
                warta_wide.generate_warta(warta_doc, prs)

    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io
