# ppt.py
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement


def format_judul_acara(text):
    if text is None:
        return ""
    text = str(text).strip()
    if "DOA PENUTUP" in text.upper():
        return "DOA PENUTUP"
    if text.endswith(":"):
        text = text[:-1].strip()

    text_upper = text.upper().replace(" ", "")

    if "BERNYANYI" in text_upper:
        pattern = r"(BERNYANYI)\s+(.*?)\s+([“\"].*[”\"])"
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            return f"{match.group(1)} {match.group(2)}\n{match.group(3)}"

    if "KOOR" in text_upper and "-" in text:
        parts = [p.strip() for p in text.split("-") if p.strip()]
        prefix_match = re.match(r"^([^:]+)", text)
        prefix = prefix_match.group(1).strip() if prefix_match else "K O O R"
        formatted_parts = []
        for p in parts:
            clean_p = re.sub(rf"^{re.escape(prefix)}", "",
                             p, flags=re.IGNORECASE).strip()
            if clean_p:
                formatted_parts.append(f"{prefix} - {clean_p}")
        return "\n".join(formatted_parts)

    return text


def apply_radical_styling(paragraph, font_size=50, alignment=PP_ALIGN.CENTER, font_color=RGBColor(0, 0, 0), font_name='Tahoma'):
    paragraph.alignment = alignment
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ['buNone', 'buAutoNum', 'buChar', 'buBlip']:
        child = pPr.find(qn(f'a:{tag}'))
        if child is not None:
            pPr.remove(child)

    buNone = OxmlElement('a:buNone')
    pPr.insert(0, buNone)
    pPr.set('marL', '0')
    pPr.set('indent', '0')

    for run in paragraph.runs:
        run.font.name = font_name
        run.font.bold = True
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color


def create_styled_slide(prs, content_text, font_size=50, alignment=PP_ALIGN.CENTER, header_text=None, font_name='Tahoma'):
    if not str(content_text).strip():
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if header_text:
        h_width = prs.slide_width - Inches(1)
        h_height = Inches(0.5)
        h_left = (prs.slide_width - h_width) // 2
        h_top = Inches(0.2)

        header_box = slide.shapes.add_textbox(h_left, h_top, h_width, h_height)
        h_p = header_box.text_frame.paragraphs[0]
        h_p.text = str(header_text).replace("\n", " ")
        apply_radical_styling(
            h_p, font_size=40, alignment=PP_ALIGN.CENTER, font_color=RGBColor(255, 165, 0), font_name='Tahoma')

    width = prs.slide_width - Inches(0.5)
    height = prs.slide_height - Inches(1.5)
    left = (prs.slide_width - width) // 2
    top = (prs.slide_height - height) // 2

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = str(content_text)
    apply_radical_styling(p, font_size=font_size,
                          alignment=alignment, font_name=font_name)


def create_cover_slide(prs, minggu, topik, tanggal):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    width_main = prs.slide_width - Inches(1)
    height_main = Inches(4)
    left_main = (prs.slide_width - width_main) // 2
    top_main = Inches(1)

    tb_main = slide.shapes.add_textbox(
        left_main, top_main, width_main, height_main)
    tf_main = tb_main.text_frame
    tf_main.word_wrap = True
    tf_main.vertical_anchor = MSO_ANCHOR.MIDDLE

    p_main = tf_main.paragraphs[0]
    p_main.text = f"{minggu}\n\"{topik}\""
    apply_radical_styling(p_main, font_size=50, font_name='Tahoma')

    width_date = prs.slide_width - Inches(1)
    height_date = Inches(1)
    left_date = (prs.slide_width - width_date) // 2
    top_date = prs.slide_height - Inches(1.5)

    tb_date = slide.shapes.add_textbox(
        left_date, top_date, width_date, height_date)
    p_date = tb_date.text_frame.paragraphs[0]
    p_date.text = tanggal
    apply_radical_styling(p_date, font_size=32, font_name='Tahoma')


def generate_slides(prs, cover_data, sections):
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    create_cover_slide(
        prs, cover_data['minggu'], cover_data['topik'], cover_data['tanggal'])

    for section in sections:
        raw_judul = section.get('judul', '').strip()
        isi_raw = section.get('isi', [])

        if not raw_judul and not isi_raw:
            continue

        raw_judul_upper = raw_judul.upper().replace(" ", "")

        is_lagu = any(k in raw_judul_upper for k in [
            "BERNYANYI", "BNSEKOLAHMINGGU", "CARIJIWABERSAMAYESUS",
            "SAPA-SAPA", "BNSM", "BESM", "KJ", "BE."
        ])

        is_epistel = "EPISTEL" in raw_judul_upper
        formatted_judul = format_judul_acara(raw_judul)

        if not is_lagu:
            create_styled_slide(prs, formatted_judul,
                                font_size=50, font_name='Tahoma')

        if isi_raw:
            if is_lagu:
                clean_lines = [line.strip() for line in isi_raw if line.strip(
                ) and not line.startswith("---")]
                for i in range(0, len(clean_lines), 2):
                    chunk = clean_lines[i:i+2]
                    content_text = "\n".join(chunk)
                    create_styled_slide(
                        prs,
                        content_text,
                        font_size=40,
                        alignment=PP_ALIGN.CENTER,
                        header_text=formatted_judul,
                        font_name='Consolas'
                    )
            else:
                for line in isi_raw:
                    clean_line = str(line).strip()
                    if clean_line and not clean_line.startswith("---"):
                        f_size = 32 if is_epistel else 36
                        create_styled_slide(
                            prs,
                            clean_line,
                            font_size=f_size,
                            alignment=PP_ALIGN.LEFT,
                            font_name='Tahoma'
                        )
